import shutil
import os
import matplotlib.pyplot as plt
import random as rand
import math
import pandas as pd
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

#@title Functions:
def create_dict(ppt_file,pointer):
    # """
    # Deletes existing images and inserts new ones in the same position.

    # :param ppt_file: Path to the input PowerPoint file.
    # :param output_file: Path to save the updated PowerPoint file.
    # :param image_replacements: Dictionary with shape names as keys and new image paths as values.
    # """
    presentation = Presentation(ppt_file)
    pointer2={}
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or True:
                try:
                  for i in pointer:
                    if str(pointer[i]) in str(shape.text).strip():
                      #print(str(shape.text))
                      if i=='50/100':
                        print(shape.text)
                      pointer2[i]=shape.name
                except(AttributeError):
                  pass
    return pointer2

def map2_creation(mapper):
  map2={}
  for i in mapper:
    if(i=="Player_img"):
      break
    map2[mapper[i]]=i
  return map2

def extract_text_format(shape):
    """
    Extracts formatting details from a shape with text.
    :param shape: A PowerPoint shape containing text.
    :return: A dictionary of text formatting properties.
    """
    if not shape.has_text_frame:
        return None

    # Get the text frame and first paragraph
    text_frame = shape.text_frame
    run=text_frame
    # paragraph = text_frame.paragraphs[0]
    # run = paragraph.runs[0] if paragraph.runs else None

    if not run:
        return None

    font = run.font
    format_details = {
        "font_name": font.name,
        "font_size": font.size.pt if font.size else None,
        "bold": font.bold,
        "italic": font.italic,
        "color": (
            font.color.rgb if font.color and font.color.type == RGBColor.TYPE_RGB else None
        ),
        "alignment": paragraph.alignment,
    }
    return format_details

def apply_text_format(shape, text, format_details):
    """
    Applies formatting details to a shape and inserts new text.
    :param shape: A PowerPoint shape where text will be written.
    :param text: Text to write into the shape.
    :param format_details: A dictionary of text formatting properties.
    """

    text_frame = shape.text_frame
    text_frame.clear()  # Clear existing text
    paragraph = text_frame.add_paragraph()
    run = paragraph.add_run()
    run.text = text

    font = run.font
    font.name = format_details.get("font_name")
    font.size = format_details.get("font_size")
    font.bold = format_details.get("bold")
    font.italic = format_details.get("italic")
    if format_details.get("color"):
        font.color.rgb = format_details["color"]
    paragraph.alignment = format_details.get("alignment")

def set_text_color(shape, color):
    """
    Sets the color of all text in a text frame to the specified RGBColor.
    :param shape: The shape containing the text frame.
    :param color: An instance of RGBColor specifying the desired color.
    """

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = color
    return shape


def update_text_split_and_format(shape, new_text):
    """
    Splits the existing text into two parts, updates the first part (number) with new_text,
    and retains the second part with its original formatting and position.
    :param shape: The shape containing the text.
    :param new_text: The new text to replace the first part of the existing text.
    """
    if not shape.has_text_frame:
        return

    # Extract the existing text
    text_frame = shape.text_frame
    paragraphs = text_frame.paragraphs

    # Ensure there is at least one paragraph and run
    if not paragraphs or not paragraphs[0].runs:
        return

    # Work with the first paragraph and its runs
    paragraph = paragraphs[0]
    runs = paragraph.runs

    # Handle cases with fewer than 2 runs
    if len(runs) < 2:
        # print("Expected two runs for formatting, handling as single word.")
        # print(new_text)
        runs[0].text = str(new_text)  # Update the single run with new text
        return

    # Get the first word (number) and second word (label)
    first_run = runs[0]  # First word (number)
    second_run = runs[1]  # Second word (label)

    # Capture formatting of the original runs
    first_format = {
        "font_name": first_run.font.name,
        "font_size": first_run.font.size,
        "bold": first_run.font.bold,
        "italic": first_run.font.italic,
        "color": first_run.font.color.rgb,
    }

    second_format = {
        "font_name": second_run.font.name,
        "font_size": second_run.font.size,
        "bold": second_run.font.bold,
        "italic": second_run.font.italic,
        "color": second_run.font.color.rgb,
    }

    # Clear the paragraph and reinsert the updated runs
    text_frame.clear()
    new_paragraph = text_frame.add_paragraph()

    # Add the updated number (first word) with preserved formatting
    new_first_run = new_paragraph.add_run()
    new_first_run.text = str(new_text)
    new_first_run.font.name = first_format["font_name"]
    new_first_run.font.size = first_format["font_size"]
    new_first_run.font.bold = first_format["bold"]
    new_first_run.font.italic = first_format["italic"]
    if first_format["color"]:
        new_first_run.font.color.rgb = first_format["color"]

    # Add the second word (label) with its original formatting
    new_second_run = new_paragraph.add_run()
    new_second_run.text = second_run.text
    new_second_run.font.name = second_format["font_name"]
    new_second_run.font.size = second_format["font_size"]
    new_second_run.font.bold = second_format["bold"]
    new_second_run.font.italic = second_format["italic"]
    if second_format["color"]:
        new_second_run.font.color.rgb = second_format["color"]

    # Updated text is now part of the shape's text frame


    # Clear the paragraph and reinsert the updated runs
    text_frame.clear()
    new_paragraph = text_frame.add_paragraph()

    # Add the updated number (first word) with preserved formatting
    new_first_run = new_paragraph.add_run()
    new_first_run.text = str(new_text)
    new_first_run.font.name = first_format["font_name"]
    new_first_run.font.size = first_format["font_size"]
    new_first_run.font.bold = first_format["bold"]
    new_first_run.font.italic = first_format["italic"]
    if first_format["color"]:
        new_first_run.font.color.rgb = first_format["color"]

    # Add the second word (label) with its original formatting
    new_second_run = new_paragraph.add_run()
    new_second_run.text = second_run.text
    new_second_run.font.name = second_format["font_name"]
    new_second_run.font.size = second_format["font_size"]
    new_second_run.font.bold = second_format["bold"]
    new_second_run.font.italic = second_format["italic"]
    if second_format["color"]:
        new_second_run.font.color.rgb = second_format["color"]

def make_colchart(slide,info,left,top,width,height):
        chart_data = CategoryChartData()
        chart_data.categories = [f'{i}' for i in range(1,11)]
        chart_data.add_series('Series 1', info)  # Example data for the series

        # Get the position of the removed image

        # Remove the old image


        # Add a bar chart in the same position
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
        ).chart

        # Customize the chart
        chart.has_title = True
        chart.chart_title.text = "Bar Chart Title"

        # Customize axes
        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = Pt(12)  # Customize category label size

        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = Pt(12)
        category_axis.tick_labels.font.color.rgb = RGBColor(255, 255, 255)
        value_axis.tick_labels.font.color.rgb = RGBColor(255, 255, 255)
        for series in chart.series:
          fill = series.format.fill
          fill.solid()
          fill.fore_color.rgb = RGBColor(255, 0, 0)
def create_filled_spider_chart(slide,info,left,top,width,height):
    """
    Creates a radar (spider) chart with a filled radar area.
    :param output_file: Path to save the PowerPoint presentation.
    """
    # Create a new PowerPoint presentation


    # Add a blank slide
    #slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank layout

    # Define radar chart data
    chart_data = CategoryChartData()
    chart_data.categories = ["Strike Rate", "Wickets", "Eco", "Matchup", "Fielding", "Avg"]
    chart_data.add_series('Player Stats', info)  # Example stats

    # Add radar chart
    #x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.RADAR_FILLED,  # Use RADAR_FILLED for filled radar chart
        left,top,width+Inches(0.5),height+Inches(0.5),chart_data
    )
    chart = chart_frame.chart

    # Customize the chart
    chart.has_title = True
    chart.chart_title.text = "Player Performance Spider Chart"

    # Set axes formatting
    category_axis = chart.category_axis
    category_axis.tick_labels.font.color.rgb = RGBColor(255,0,0)  # White category labels

    value_axis = chart.value_axis
    value_axis.maximum_scale = 150  # Adjust max scale based on your stats
    value_axis.minimum_scale = -50
    value_axis.has_major_gridlines = True
    value_axis.visible = False
    category_axis.visible=False
    # change_gridline_color(value_axis, (255,255,255))
    # change_gridline_color(category_axis, (255,255,255))
    category_axis.tick_labels.font.size = Pt(15)

    # Customize the data series (radar area)
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = RGBColor(255, 100, 100)  # Light red fill
def get_picto(ppt_path):
# Load the PowerPoint presentation
  presentation = Presentation(ppt_path)

  # Load the PowerPoint presentation
  rico={}

  # Loop through slides and shapes to find images with alt text
  for slide_number, slide in enumerate(presentation.slides, start=1):
      for shape_number, shape in enumerate(slide.shapes, start=1):
          if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # Check if the shape is a picture
              # Access the alt text via the XML
              alt_text = shape._element.xpath(".//p:cNvPr/@descr")
              if alt_text:  # Check if there's alt text
                  try:
                    print(picto['alt_text[0]'])
                  except:
                    pass
                  print(f"Slide {slide_number}, Shape {shape.name}: Alt text: {alt_text[0]}")
                  rico[alt_text[0]]=shape.name
              else:
                  print(f"Slide {slide_number}, Shape {shape_number}: No alt text")
  return rico
#@title graph_functions


def create_dual_area(matches,OPI,width,height,baseline=1.0):
    """
    Creates a dual-area chart with diagonal lines and filled regions.
    :param x: List of x-coordinates.
    :param y: List of y-coordinates for the graph.
    :param baseline: The baseline value.
    :param output_image: Path to save the chart image.
    """
    # Plot setup
    #opi is list of all 10
    #matches is x axis (0->10)
    width/=914400
    height/=914400
    plt.figure(figsize=(width, height))
    x=matches
    y=OPI

    # Plot the chart
    plt.plot(x,y, color='purple', linewidth=2, label='Graph', marker='o')  # Line with markers
    plt.axhline(y=baseline, color='blue', linestyle='--', label='Baseline')  # Baseline

    # Fill areas between graph and baseline
    for i in range(len(x) - 1):
        # Segment coordinates
        x_segment = [x[i], x[i + 1]]
        y_segment = [y[i], y[i + 1]]

        # Fill green if above baseline
        if all(yi > baseline for yi in y_segment):
            plt.fill_between(x_segment, baseline, y_segment, color='green', alpha=0.8)
        # Fill red if below baseline
        elif all(yi <= baseline for yi in y_segment):
            plt.fill_between(x_segment, baseline, y_segment, color='red', alpha=0.8)
        # Mixed segment
        else:
            intersection_x = x[i] + (baseline - y[i]) * (x[i + 1] - x[i]) / (y[i + 1] - y[i])
            if y[i] > baseline:
                plt.fill_between([x[i], intersection_x], baseline, [y[i], baseline], color='green', alpha=0.8)
                plt.fill_between([intersection_x, x[i + 1]], baseline, [baseline, y[i + 1]], color='red', alpha=0.8)
            else:
                plt.fill_between([x[i], intersection_x], baseline, [y[i], baseline], color='red', alpha=0.8)
                plt.fill_between([intersection_x, x[i + 1]], baseline, [baseline, y[i + 1]], color='green', alpha=0.8)
    #plt.axis('off')
    plt.gca().set_facecolor('black')
    plt.xlabel('X-axis', color='white')
    plt.ylabel('Y-axis', color='white')
    plt.title('Dual Area Chart with Diagonal Lines', color='white')
    plt.legend(facecolor='black', edgecolor='white')  # Legend background and border
    # Customize axis spines (make them visible on black background)
    plt.gca().spines['bottom'].set_color('white')  # X-axis line color
    #plt.gca().spines['bottom'].set_linewidth(0.5)  # X-axis line thickness

    plt.gca().spines['left'].set_color('black')  # Y-axis line color
    #plt.gca().spines['left'].set_linewidth(0.5)  # Y-axis line thickness

    plt.gca().spines['top'].set_color('none')  # Hide top spine
    plt.gca().spines['right'].set_color('none')  # Hide right spine

    # Labels and legend
    plt.xlabel('X-axis')
    plt.ylabel('Y-axis')
    plt.title('Dual Area Chart with Diagonal Lines')
    plt.legend()

    # Save the chart as an image
    base_dir = Path(os.path.dirname(os.path.abspath(__file__)))
    plt.savefig(f'{base_dir}/pptx/file.png', bbox_inches='tight', pad_inches=0, dpi=300)

    plt.close()



# Example data (discrete)


# Create and insert the chart


def create_donut_chart(consistent):
    """
    Creates a donut chart with percentages and labels.
    :param wins: List of wins for each category.
    :param matches: List of matches for each category.
    :param output_file: Path to save the chart as an image.
    """
    base_dir = Path(os.path.dirname(os.path.abspath(__file__)))
    output_file=f'{base_dir}/pptx/file.png'
    # Calculate percentages
    percentages = [(i) * 100 for i in consistent]
    labels = ['Overall', 'vs Team', 'Venue']  # Example labels for each category

    # Chart setup
    fig, axs = plt.subplots(1, 3, figsize=(12, 6), subplot_kw={'aspect': 'equal'})
    fig.patch.set_facecolor('black')  # Set background color to black

    # Colors
    primary_color = '#ff4c4c'  # Red for the main segment
    secondary_color = '#3d3d3d'  # Dark gray for the remaining segment
    text_color = 'white'  # White for text

    for i, ax in enumerate(axs):
        # Pie chart for each category
        sizes = [percentages[i], 100 - percentages[i]]
        colors = [primary_color, secondary_color]
        wedges, texts = ax.pie(
            sizes,
            colors=colors,
            startangle=90,
            wedgeprops={'width': 0.3, 'edgecolor': 'black'},
        )

        # Add the percentage in the center
        ax.text(0, 0, f"{int(percentages[i])}%", ha='center', va='center', fontsize=18, color=text_color, fontweight='bold')

        # Add the label below the donut chart
        ax.text(0, -1.4, labels[i], ha='center', va='center', fontsize=14, color=text_color, fontweight='bold')

    # Adjust layout
    plt.subplots_adjust(wspace=0.4)
    plt.savefig(output_file, bbox_inches='tight', dpi=300, facecolor=fig.get_facecolor())
    plt.close()
    print(f"Donut chart saved as: {output_file}")
#@title fill_text
def fill_text(template_file, output_file, mapping,replace,map3,n):
    """
    Automates filling a PowerPoint template with new images and text.
    A new copy of the template is created each time to preserve the original.

    :param template_file: Path to the original PowerPoint template file
    :param output_file: Path to save the updated PowerPoint file
    :param replacements: Dictionary with keys as placeholder names and values as
                         the new content (text or image path)
    """
    # Create a copy of the template to work on
    base_dir = Path(os.path.dirname(os.path.abspath(__file__)))
    temp_dir = base_dir / "pptx"
    temp_dir.mkdir(exist_ok=True)
    destination_folder = temp_dir
    shutil.copy(template_file, destination_folder)
    temp_copy = f"{destination_folder}/batting{n}.pptx"
    # Load the copied PowerPoint file
    presentation = Presentation(temp_copy)
    #team_name="I C U"
    for slide in presentation.slides:
        for shape in slide.shapes:
            # Check for the specific placeholder for team_name
            if shape.name == "Google Shape;37;p1":  # Specific name to handle
                pass
                #nothing happens here

            # Replace images for other shapes
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if shape.name in map3:  # Use shape name as a placeholder
                    if map3[shape.name]=='r_player':
                      base_dir = Path(os.path.dirname(os.path.abspath(__file__)))
                      new_image_path = f'{base_dir}/{replace["player_img"]}'#IMG
                      # Get the current position and size of the image
                      left = shape.left
                      top = shape.top
                      width = shape.width
                      height = shape.height

                      # Remove the old image
                      sp = shape._element
                      sp.getparent().remove(sp)

                      # Add the new image
                      slide.shapes.add_picture(new_image_path, left, top, width=width/1.5, height=height/1.5)
                    if map3[shape.name]=='team_symbol':
                      base_dir = Path(os.path.dirname(os.path.abspath(__file__)))
                      new_image_path =f'{base_dir}/{replace["team_symbol"]}'#IMG

                      left = shape.left
                      top = shape.top
                      width = shape.width
                      height = shape.height

                      # Remove the old image
                      sp = shape._element
                      sp.getparent().remove(sp)

                      # Add the new image
                      slide.shapes.add_picture(new_image_path, left, top, width=width, height=height)
                    if map3[shape.name]=='spidergraph':
                      left = shape.left
                      top = shape.top
                      width = shape.width
                      height = shape.height

                      # Remove the old image
                      sp = shape._element
                      sp.getparent().remove(sp)

                      rect = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,  # Shape type: Rounded rectangle
                            left/1.03, top*1.025, width*1.15, height      # Position and size
                        )
                      #["Strike Rate", "Wickets", "Eco", "Matchup", "Fielding", "Avg"]
                      info=[replacements['strike_rate'],replacements['wickets'],replacements['economy'],replacements['form'],replacements['catches'],replacements['Average']]
                      # Set the fill color to black
                      fill = rect.fill
                      fill.solid()
                      fill.fore_color.rgb = RGBColor(50,50,50)
                      create_filled_spider_chart(slide,info,left,top,width,height)
                      pass
                    if map3[shape.name]=='graph':
                      if n==1:
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height

                        # Remove the old image
                        sp = shape._element
                        sp.getparent().remove(sp)
                        rounded_rect = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,  # Shape type: Rounded rectangle
                            left, top, width, height     # Position and size
                        )

                        # Set the fill color to black
                        fill = rounded_rect.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(0, 0, 0)
                        #slide.shapes.add_picture(new_image_path, left, top, width=width/1.5, height=height/1.5)
                        info= replacements['graph']['fantasy']
                        print(info)
                        make_colchart(slide,info,left,top,width,height)

                      if n==2:

                        consistent=replacements['graph']['consistancy']
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height

                        rounded_rect = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,  # Shape type: Rounded rectangle
                            left, top, width, height*1.01     # Position and size
                        )

                        # Set the fill color to black
                        fill = rounded_rect.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(0, 0, 0)

                        create_donut_chart(consistent)
                        base_dir = Path(os.path.dirname(os.path.abspath(__file__)))
                        new_image_path = f'{base_dir}/pptx/file.png'
                        slide.shapes.add_picture(new_image_path, left, top, width=width, height=height)
                      if n==3:

                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        sp = shape._element
                        sp.getparent().remove(sp)
                        rounded_rect = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,  # Shape type: Rounded rectangle
                            left, top, width, height     # Position and size
                        )
                        vpi=replacements['graph']['vpi']
                        matches=[i for i in range(1,len(vpi)+1)]

                        print(width,height)
                        create_dual_area(matches,vpi,width,height,baseline=1.0)
                        # Set the fill color to black
                        fill = rounded_rect.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(0, 0, 0)
                        # Remove the old image
                        base_dir = Path(os.path.dirname(os.path.abspath(__file__)))
                        new_image_path = f'{base_dir}/pptx/file.png'
                        slide.shapes.add_picture(new_image_path, left, top, width=width, height=height)
                      if n==4:
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        sp = shape._element
                        sp.getparent().remove(sp)
                        rounded_rect = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,  # Shape type: Rounded rectangle
                            left, top, width, height     # Position and size
                        )
                        opi=replacements['graph']['opi']
                        matches=[i for i in range(1,len(opi)+1)]
                        print(width,height)
                        create_dual_area(matches,opi,width,height,baseline=1.0)
                        # Set the fill color to black
                        fill = rounded_rect.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(0, 0, 0)
                        # Remove the old image
                        base_dir = Path(os.path.dirname(os.path.abspath(__file__)))
                        new_image_path = f'{base_dir}/pptx/file.png'
                        slide.shapes.add_picture(new_image_path, left, top, width=width, height=height)

    for slide in presentation.slides:
        for shape in slide.shapes:
            # Replace text in text boxes
            if shape.name in mapping:

              new_text = replacements[mapping[shape.name]]

              update_text_split_and_format(shape, new_text)
              #print(run.font.color.rgb)


    # Save the updated PowerPoint


    presentation.save(output_file)

    # Remove the temporary copy of the template
    os.remove(temp_copy)

    print(f"Updated presentation saved as: {output_file}")

#@title (batting) pointer
pointer = {
    "name":"VIRAT",
    "team":"ROYAL",
    "expt_pts": 87,
    "cost": 11,
    "expt_runs": 99,
    "expt_4s_and_6s":'7-',

    "runs": 49,#both expt_runs and runs store same val
    "strike_rate": 220,
    "boundaries": '15-',
    "wickets": '2-',

    "economy": 12,
    "catches": '13-',
    "run_outs": '1-',
    "form": 92,

    "Matches": 244,
    "cruns":42,
    "hstrike_rate":154,
    "4s/6s":'44',#44/76 but 44 is enough to find it

    "highest_score": 102,
    "innings": 181,
    "Average":53,
    "50/100":"21",#21/25 actually
}

#@title map2 creation
def map2_creation(mapper):
  map2={}
  for i in mapper:
    if(i=="Player_img"):
      break
    map2[mapper[i]]=i
  return map2
#@title map3 creation
def map3_creation(picto):
  map3={}
  map3[picto["graph"]]='graph'
  map3[picto["spiderman"]]='spidergraph'
  map3[picto["r_player"]]='r_player'
  map3[picto["team_symbol"]]='team_symbol'
  return map3
def PPT_GEN(n,replacements,pointer):#n is slide number
  base_dir = Path(os.path.dirname(os.path.abspath(__file__)))
  template_dir = base_dir / "pptx_template"
  output_dir = base_dir / "pptx_out"
  template_dir.mkdir(exist_ok=True)
  output_dir.mkdir(exist_ok=True)

  in_file = template_dir / f"batting{n}.pptx"
  out_file = output_dir / f"ppt{n}.pptx"

  if not in_file.exists():
    raise FileNotFoundError(f"Template file not found at {in_file}")


  mapper=create_dict(in_file, pointer)

  picto=get_picto(in_file)
  #print(picto)
  mapper["Player_img"]=picto["r_player"] #if there is error here you may have forgotten to import pickle from binary
  mapper["graph"]=picto["graph"]
  mapper["spiderman"]=picto["spiderman"]
  mapper["team_symbol"]=picto["team_symbol"]
  #mapper["team_name"]=picto["team_name"]
  #mapper["team_symbol"]=picto["team_symbol"]
  print(mapper)
  #mapper completion
  map3=map3_creation(picto)
  map2=map2_creation(mapper)
  fill_text(in_file, out_file,map2,replacements,map3,n)#inputfile,output file, mapping(shape.name->human_name),replacements(human_name->value)
  #replace_images(mapper,image_file)
#####################################
#@title example replacement data
#@title replacements
replacements = {
    "name":"Bob ",
    "player_id":"VIRAT001",
    "team":"Cartoon Network",
    "expt_pts": 32,
    "cost": 15,
    "expt_runs": 77,
    "expt_4s_and_6s": 8,
    "runs": 77,#both expt_runs and runs store same val
    "strike_rate": 8,
    "boundaries": 8,
    "wickets": 1,
    "economy": 57,
    "catches": 23,#fielding stats->catches
    "run_outs": 3,
    "form": 92,
############################### upto here, independent of n
    "Matches": 244,
    "cruns":324089,
    "hstrike_rate":23,
    "4s/6s":'134/45',
    "player_img":'image.jpg',#hardcoded file path into fill_text for now
    "team_symbol":'image.jpg',#no code written to replace symbol currently. shape.name is automatically located

    "highest_score": 55,#
    "innings": 3489,
    "Average":78,
    "50/100":'133/44',
    'graph':{
        'fantasy':[1,2,3,4,5,6,7,8,9,10],#[last 10 days of fantasy points]when (n==1)
        'consistancy':[0.5,0.75,0.9],#[overall,against given opp, in given venue]when (n==2)
        'vpi':[1,1.5,0.5,0.4,0.25,6,0.75,1.8,0.1,1.9,0.9],#[last 10 vpi]-------- (n==3)[if current venue data is not possible, skip venue]
        'opi':[1,1.5,0.5,0.4,0.25,6,0.75,1.8,0.1,1.9,0.9],#[last 10 opi] when (only n=4)
    },
    'n':1

}
#n==1: overall
#n==2: past 10 days
#n==3: in given venue
#n==4: against given opponent
    # "Matches": 244,
    # "cruns":324089,
    # "hstrike_rate":23,
    # "4s/6s":'134/45',
    # "highest_score": 55,
    # "innings": 3489,
    # "Average":78,
    # "50/100":'133/44'
#search IMG in fill text to find  player and symbol image path code
for i in range (1, 5):
    PPT_GEN(i,replacements,pointer)
#output_file='/content/file.png' here i store the graphs temporarily. have not written deletion code.(note this file is created automatically by my code)
